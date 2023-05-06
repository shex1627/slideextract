from PIL import Image
import os
import argparse
from typing import List, Dict
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import re
from datetime import datetime, timedelta
import boto3
from io import BytesIO
from pptx.util import Emu

def numerical_to_datetime(number: str) -> str:
    """
    Converts a numerical string to a timestamp string in the format "HH:MM:SS".

    Args:
        number (str): A numerical string representing a timestamp.

    Returns:
        str: A string representing the timestamp in the format "HH:MM:SS".
    """
    timestamp = timedelta(seconds=int(number))
    return str(timestamp)

def create_pptx(images_dict: Dict[int, Image.Image]) -> Presentation:
    """
    Creates a new PowerPoint presentation with the given list of images.

    Args:
        images (list): A list of PIL image objects to be added to the presentation.

    Returns:
        Presentation: A new PowerPoint presentation object with the images added to slides.

    Raises:
        ValueError: If any of the images are not in RGB mode.
    """
    # Set up the presentation with default slide size (10 inches x 7.5 inches)
    prs = Presentation()
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches

    # Add each image to a new slide
    for image_index, image in images_dict.items():
        # Check if image is in RGB mode
        if image.mode != 'RGB':
            raise ValueError(f"Image is not in RGB mode.")

        # Check if image file name is a number
        match = re.match(r"^\d+$", str(image_index))
        if match:
            # Convert the numerical file name to a datetime object and add a note
            timestamp = numerical_to_datetime(match.group())
            note = f"Note: This image was taken at {timestamp}."
        else:
            # No note is required for non-numerical file names
            note = ""
        
        # Add the image and note to a new slide
        buffer = BytesIO()
        image.save(buffer, format="jpeg")
        buffer.seek(0)
        
        # Resize the image to fit the slide dimensions
        width, height = image.size
        prs.slide_width = int(width * (914400/96))
        prs.slide_height = int(height * (914400/96))
        # add slide to ppt
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left = top = 0
        pic = slide.shapes.add_picture(buffer, left, top, width=prs.slide_width, height=prs.slide_height)
        if note:
            slide.notes_slide.notes_text_frame.text = note
    
    return prs


def upload_pptx_to_s3(pptx: Presentation, s3_bucket: str, s3_folder: str, title: str) -> str:
    """
    Uploads the given PowerPoint presentation to the specified S3 bucket and folder.

    Args:
        pptx (Presentation): The PowerPoint presentation to upload.
        s3_bucket (str): The name of the S3 bucket to upload the presentation to.
        s3_folder (str): The name of the S3 folder to upload the presentation to.
        title (str): The title of the PowerPoint presentation.

    Returns:
        str: The S3 URL of the uploaded presentation.
    """
    # Convert the PowerPoint presentation to a bytes buffer
    pptx_bytes = BytesIO()
    pptx.save(pptx_bytes)
    pptx_bytes.seek(0)

    # Upload the bytes buffer to the specified S3 bucket and folder
    s3 = boto3.resource('s3')
    key = f"{s3_folder}/{title}.pptx"
    s3.Bucket(s3_bucket).put_object(Key=key, Body=pptx_bytes)

    # Return the S3 URL of the uploaded presentation
    return f"https://{s3_bucket}.s3.amazonaws.com/{key}"
